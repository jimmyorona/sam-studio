#!/usr/bin/env python3
"""Multi-persona document review / rewrite, bridged into the SAM Studio web app.

A one-shot CLI shaped like pptx_to_video.py: Express spawns it as a subprocess,
captures stdout for the SSE log stream, and reads the markdown reports it writes
to reviews/<doc-slug>/.

Ported from sam-slide-reviewer/webapp/server.py (the retired FastAPI app):
carries over REVIEWER/REWRITER/SYNTHESIS prompts, the concurrent fan-out
(MAX_CONCURRENT_REVIEWS), the synthesis pass, and the DOCX / PPTX export helpers.
Uses `requests` + a thread pool instead of httpx/asyncio so it has the same
dependency surface as the rest of scripts/.

Subcommands
-----------
  run     extract -> fan-out persona reviews -> synthesis, writing reports
  export  render a finished report markdown file to .docx or .pptx

Progress lines are printed to stdout with an [HH:MM:SS] prefix (human log) and,
for machine parsing by the server, structured markers:

  @@STATE persona=<slug> state=<queued|running|done|error>
  @@SLUG <doc-slug>
  @@REPORT slug=<file-stem> name=<display name>
  @@DONE state=<complete|error> error=<msg>

Env:
  OLLAMA_URL  (default http://localhost:11434; --ollama-url overrides)
"""

import argparse
import datetime
import io
import os
import re
import sys
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import requests

ROOT = Path(__file__).resolve().parent.parent
PERSONAS_DIR = ROOT / "personas"
REVIEWS_DIR = ROOT / "reviews"
DEFAULT_OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
MAX_CONCURRENT_REVIEWS = int(os.environ.get("REVIEW_CONCURRENCY", "3"))
SUPPORTED_EXTS = {".pptx", ".pdf", ".docx", ".doc", ".odt", ".md", ".markdown", ".txt"}

sys.path.insert(0, str(ROOT / "scripts"))
import extract  # noqa: E402  (reuses the CLI extractors)

REVIEWER_SYSTEM_PROMPT = """\
You are an expert document reviewer. Your reviewing identity is defined by the
persona brief included in the user message. Fully adopt that persona's
archetype, values, standards, tone, and priorities.

How to interpret the persona brief:
- AUTHORITATIVE: the "Review Lens (Document & Slide Review)" section. Follow
  its writing style, hunt priorities, and severity calibration exactly, and
  use its scorecard dimensions verbatim.
- ADOPT: the Character Profile, Voice Personality, and guardrails ("What SAM
  Never Does").
- IGNORE: all TTS mechanics (ElevenLabs settings, SSML, voice recommendations).

Write the review in this exact markdown structure:

# <Document name> — Review by <Persona name>

> One-paragraph overall verdict, in the persona's voice.

## Strengths
(2–5 bullets)

## Findings
### <N>. <Short title> — `Slide X` / `Page X` / `Section`
**Severity:** Critical / Major / Minor
**Issue:** ...
**Recommendation:** ... (with rewritten text where useful)

## What this persona would ask the author
(2–4 pointed questions)

## Scorecard
| Dimension | Score (1–5) | Note |
(Use the persona's Review Lens scorecard dimensions.)

Rules: reference exact slide/page numbers from the extracted content headings;
every finding must be concrete and actionable; if the persona's domain barely
applies to this document, say so and keep the review short. Output only the
markdown report — no preamble.
"""

REWRITER_SYSTEM_PROMPT = """\
You are an expert presentation writer. Your writing identity is defined by the
persona brief included in the user message. Fully adopt that persona's
archetype, tone, structure preferences, and standards.

How to interpret the persona brief:
- The "Hybrid Tone Rules" table defines the texture to use per section type.
- "What SAM Never Does" entries are hard constraints on your writing.
- The "Review Lens (Document & Slide Review)" section defines the quality bar:
  the rewrite must be a presentation that THIS persona's own review would score
  5/5 on every scorecard dimension.
- IGNORE all TTS mechanics (ElevenLabs settings, SSML, voice recommendations).

Task: rewrite the ENTIRE presentation, slide by slide, in this persona's voice
and structure.

Rules:
- Keep the `## Slide N` heading structure, one section per slide:
  `## Slide N: <title>`.
- You may reorder, merge, or split slides when the persona's structure demands
  it (e.g. conclusion-first). When content moves, append the source in the
  heading: `## Slide 1: <title> (source: slide 5)`.
- Preserve every fact and figure exactly. NEVER invent data, numbers, sources,
  or claims. Where the persona's standards require something the source lacks
  (a baseline, a source citation, an explicit ask), insert a visible
  placeholder instead: `[NEEDS: baseline figure for the $2M claim]`.
- For each slide provide the on-slide content (concise bullets or short lines,
  shaped how this persona would shape them), then `**Speaker notes:**` —
  the spoken narration in the persona's full voice.
- If the user message includes a "REVIEW FINDINGS TO APPLY" section, treat it as
  a required work order: resolve every priority fix, consensus finding, and
  flagged issue it lists, in addition to your own persona's standards. These
  findings come from a prior expert review of THIS document and take precedence
  over your persona's stylistic preferences where they conflict. Where a finding
  asks for data the source lacks, insert a `[NEEDS: …]` placeholder rather than
  inventing it.
- After the last slide, add a `## Rewrite notes` section: 3–6 bullets on the
  structural changes made and why, plus a list of all `[NEEDS: …]` placeholders.
  When review findings were supplied, briefly note how each priority fix was
  addressed.

Output only the rewritten presentation in markdown — no preamble.
"""

REWRITER_ADVISE_ADDENDUM = """\

ADVISE MODE (enabled): do not leave gaps blank. For every `[NEEDS: X]`
placeholder you would write, immediately follow it with `[DRAFT: …]` containing
a specific, plausible value or sentence a domain expert might reasonably expect,
so the author can confirm or correct it instead of starting from nothing. Base
each draft on the document's context, the review findings, and conventional
benchmarks — but never present a draft as verified fact; it is a starting
proposal. Always keep the `[NEEDS: …]` marker in place so the gap stays tracked,
with the `[DRAFT: …]` right after it.
"""

SYNTHESIS_SYSTEM_PROMPT = """\
You are the chair of an expert review panel. You will receive several reviews
of the same document, each written by a different expert persona. Write a
synthesis in this exact markdown structure:

# <Document name> — Review Synthesis
*Panel: <persona names> · <date>*

## Overall verdict
(2–3 sentences)

## Consensus findings
(Issues flagged by 2+ personas; for each: the issue, who flagged it, the fix.)

## Conflicts & tensions
(Where personas disagree; name the trade-off and recommend a resolution. Each
review ends with that persona's known blind spots — use them to weigh sides.)

## Unique catches
(Important findings only one persona caught.)

## Top 5 priority fixes
(Ordered, actionable, with slide/page references.)

Output only the markdown report — no preamble.
"""


# --------------------------------------------------------------------------- #
# Logging / progress markers (parsed by web/server/index.js)
# --------------------------------------------------------------------------- #
def log(msg: str) -> None:
    print(f"[{time.strftime('%H:%M:%S')}] {msg}", flush=True)


def marker(line: str) -> None:
    print(line, flush=True)


# --------------------------------------------------------------------------- #
# Persona helpers (ported from server.py)
# --------------------------------------------------------------------------- #
def persona_display_name(path: Path) -> str:
    stem = re.sub(r"^\d+-", "", path.stem)
    stem = re.sub(r"-PERSONA$", "", stem)
    return stem.replace("-", " ").title()


def doc_slug(filename: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", Path(filename).stem.lower()).strip("-") or "document"


def extract_content(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pptx":
        body = extract.extract_pptx(path)
    elif ext == ".pdf":
        body = extract.extract_pdf(path)
    elif ext in (".docx", ".doc", ".odt"):
        body = extract.extract_docx(path)
    else:
        body = path.read_text(encoding="utf-8", errors="replace")
    return f"# Extracted content: {path.name}\n\n{body}"


def gather_review_findings(out_dir: Path) -> str:
    """Collect prior review output for this document so the rewrite can apply it.

    Prefers the consolidated synthesis; falls back to concatenating the
    individual persona reviews. Rewrite outputs and the raw extract are skipped.
    Returns "" when no prior review exists.
    """
    synth = out_dir / "00-SYNTHESIS.md"
    if synth.exists():
        return synth.read_text(encoding="utf-8").strip()
    parts = []
    for f in sorted(out_dir.glob("*.md")):
        if f.name == "_extracted.md" or f.name == "00-SYNTHESIS.md" or f.name.startswith("rewrite-"):
            continue
        parts.append(f"=== REVIEW: {f.stem} ===\n\n{f.read_text(encoding='utf-8')[:6000]}")
    return "\n\n".join(parts).strip()


_CTX_CACHE: dict[str, "int | None"] = {}


def model_max_ctx(ollama_url: str, model: str) -> "int | None":
    """The selected model's maximum context window, via Ollama /api/show.

    Set REVIEWER_NUM_CTX to override (e.g. to cap memory use on a small GPU).
    Cached per model. Returns None if it can't be determined (then Ollama's
    own default applies).
    """
    override = os.environ.get("REVIEWER_NUM_CTX")
    if override and override.isdigit():
        return int(override)
    if model in _CTX_CACHE:
        return _CTX_CACHE[model]
    n = None
    try:
        resp = requests.post(f"{ollama_url}/api/show", json={"model": model}, timeout=(10, 30))
        resp.raise_for_status()
        info = resp.json().get("model_info", {}) or {}
        for key, val in info.items():
            if key.endswith(".context_length") and isinstance(val, int):
                n = val
                break
    except Exception:  # noqa: BLE001  (best-effort; fall back to Ollama default)
        n = None
    _CTX_CACHE[model] = n
    return n


def ollama_chat(ollama_url: str, model: str, system: str, user: str,
                num_ctx: "int | None" = None) -> str:
    payload = {
        "model": model,
        "stream": False,
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
    }
    if num_ctx:
        payload["options"] = {"num_ctx": num_ctx}
    resp = requests.post(f"{ollama_url}/api/chat", json=payload, timeout=(10, 900))
    resp.raise_for_status()
    return resp.json()["message"]["content"].strip()


def clear_stale_outputs(out_dir: Path, rewrite: bool) -> None:
    """Remove this mode's previous outputs so each run starts clean.

    A review run clears prior persona reviews + synthesis; a rewrite run clears
    prior rewrites. The raw extract and the *other* mode's files are kept — a
    rewrite reads the review's reports for its findings, so a rewrite must not
    delete them (and vice-versa).
    """
    removed = 0
    for f in out_dir.glob("*.md"):
        if f.name == "_extracted.md":
            continue
        is_rewrite = f.name.startswith("rewrite-")
        if is_rewrite == rewrite:
            f.unlink()
            removed += 1
    if removed:
        log(f"  Cleared {removed} stale {'rewrite' if rewrite else 'review'} report(s) from this slug.")


# --------------------------------------------------------------------------- #
# run: fan-out reviews + synthesis
# --------------------------------------------------------------------------- #
def run(args) -> int:
    rewrite = args.mode == "rewrite"
    out_dir = REVIEWS_DIR / args.slug
    out_dir.mkdir(parents=True, exist_ok=True)

    # Resolve extracted content: from --extracted, --input file, or --text.
    if args.input:
        log(f"Extracting {Path(args.input).name} …")
        try:
            content = extract_content(Path(args.input))
        except Exception as e:  # noqa: BLE001
            marker(f"@@DONE state=error error=Extraction failed: {e}")
            log(f"ERROR: extraction failed: {e}")
            return 1
    elif args.text:
        content = f"# Extracted content: {args.file_name}\n\n{args.text.strip()}"
    else:
        marker("@@DONE state=error error=Provide --input or --text")
        return 2

    if len(content.strip()) < 80:
        marker("@@DONE state=error error=Almost no text to work with (image-only or empty)")
        log("ERROR: almost no extractable text")
        return 1

    (out_dir / "_extracted.md").write_text(content, encoding="utf-8")
    marker(f"@@SLUG {args.slug}")

    # A fresh review replaces any prior review of this document (rewrites are
    # cleared later, after their findings are read).
    if not rewrite:
        clear_stale_outputs(out_dir, rewrite=False)

    persona_files = [f for f in args.personas.split(",") if f.strip()]
    personas = [
        {"file": f, "slug": Path(f).stem.lower(),
         "name": persona_display_name(Path(f)), "state": "queued", "error": None}
        for f in persona_files
    ]
    for p in personas:
        marker(f"@@STATE persona={p['slug']} state=queued")

    verb = "rewrite" if rewrite else "review"
    log(f"Running {len(personas)} persona {verb}(s), up to {MAX_CONCURRENT_REVIEWS} concurrent …")

    # Use the model's full context window so stacked prompts (persona brief +
    # context + findings + document) are never silently truncated.
    num_ctx = model_max_ctx(args.ollama_url, args.model)
    if num_ctx:
        log(f"  Model context window: {num_ctx} tokens.")

    # Carry the prior review's analysis into every rewrite, regardless of the
    # rewrite persona. Gather once up front so concurrent rewrites don't read
    # each other's partial output.
    findings_block = ""
    if rewrite:
        findings = gather_review_findings(out_dir)
        if findings:
            findings_block = (
                "REVIEW FINDINGS TO APPLY (from a prior multi-persona review of this "
                "document — address these regardless of your own persona's lens):\n\n"
                f"{findings}\n\n---\n\n"
            )
            log(f"  Applying prior review findings to the rewrite ({len(findings)} chars).")
        else:
            log("  No prior review found for this document — rewriting from the source only.")
        # Findings are captured; now drop any prior rewrites for this slug.
        clear_stale_outputs(out_dir, rewrite=True)

    rewriter_system = REWRITER_SYSTEM_PROMPT
    if rewrite and args.advise:
        rewriter_system += REWRITER_ADVISE_ADDENDUM
        log("  Advise mode on — drafting proposed content for identified needs.")

    # Optional background context (typed text and/or an attached file).
    context_text = args.context or ""
    if args.context_file:
        try:
            context_text = (context_text + "\n\n" + Path(args.context_file).read_text(
                encoding="utf-8", errors="replace")).strip()
        except OSError as e:
            log(f"  WARNING: could not read context file: {e}")
    context_block = ""
    if context_text.strip():
        context_block = (
            "BACKGROUND CONTEXT (provided by the user to inform your "
            f"{'rewrite' if rewrite else 'review'} — use it for grounding, but the "
            "document content below is what you assess):\n\n"
            f"{context_text.strip()}\n\n---\n\n"
        )
        log(f"  Using {len(context_text)} chars of background context.")

    def review_one(p: dict) -> dict:
        p["state"] = "running"
        marker(f"@@STATE persona={p['slug']} state=running")
        log(f"  {p['name']}: {verb} started")
        try:
            brief = (PERSONAS_DIR / p["file"]).read_text(encoding="utf-8")
            if rewrite:
                system = rewriter_system
                task_line = f"PRESENTATION TO REWRITE ({args.file_name})"
                out_file = f"rewrite-{p['slug']}.md"
            else:
                system = REVIEWER_SYSTEM_PROMPT
                task_line = f"DOCUMENT TO REVIEW ({args.file_name})"
                out_file = f"{p['slug']}.md"
            user_msg = (
                f"PERSONA BRIEF (referred to as \"{p['name']}\" — use that name "
                f"where a persona name is needed, not the internal character name):\n\n"
                f"{brief}\n\n---\n\n"
                f"{context_block}"
                f"{findings_block if rewrite else ''}"
                f"{task_line}:\n\n{content}"
            )
            report = ollama_chat(args.ollama_url, args.model, system, user_msg, num_ctx)
            (out_dir / out_file).write_text(report + "\n", encoding="utf-8")
            p["state"] = "done"
            marker(f"@@STATE persona={p['slug']} state=done")
            marker(f"@@REPORT slug={Path(out_file).stem} name={p['name']}"
                   + (" Rewrite" if rewrite else ""))
            log(f"  {p['name']}: done")
        except Exception as e:  # noqa: BLE001  (surface per-persona failures)
            p["state"] = "error"
            p["error"] = str(e)
            marker(f"@@STATE persona={p['slug']} state=error")
            log(f"  {p['name']}: ERROR {e}")
        return p

    with ThreadPoolExecutor(max_workers=MAX_CONCURRENT_REVIEWS) as pool:
        futures = [pool.submit(review_one, p) for p in personas]
        for _ in as_completed(futures):
            pass

    done = [p for p in personas if p["state"] == "done"]
    if not done:
        err = personas[0].get("error") or "unknown error"
        marker(f"@@DONE state=error error=All persona {verb}s failed: {err}")
        log("ERROR: all personas failed")
        return 1

    if rewrite:
        marker("@@DONE state=complete")
        log("Rewrite complete.")
        return 0

    if len(done) < 2:
        log("Single reviewer — nothing to merge, skipping synthesis.")
        marker("@@DONE state=complete")
        return 0

    log("Synthesizing panel reviews …")
    try:
        parts = []
        for p in done:
            text = (out_dir / f"{p['slug']}.md").read_text(encoding="utf-8")
            parts.append(f"=== REVIEW BY {p['name'].upper()} ===\n\n{text[:8000]}")
        user_msg = (
            f"Document: {args.file_name}\n"
            f"Date: {datetime.date.today().isoformat()}\n"
            f"Panel: {', '.join(p['name'] for p in done)}\n\n"
            + "\n\n".join(parts)
        )
        synthesis = ollama_chat(args.ollama_url, args.model, SYNTHESIS_SYSTEM_PROMPT, user_msg, num_ctx)
        (out_dir / "00-SYNTHESIS.md").write_text(synthesis + "\n", encoding="utf-8")
        marker("@@REPORT slug=00-SYNTHESIS name=Synthesis")
        marker("@@DONE state=complete")
        log("Synthesis complete.")
        return 0
    except Exception as e:  # noqa: BLE001
        marker(f"@@DONE state=error error=Synthesis failed: {e}")
        log(f"ERROR: synthesis failed: {e}")
        return 1


# --------------------------------------------------------------------------- #
# export: report markdown -> DOCX / PPTX (ported from server.py)
# --------------------------------------------------------------------------- #
def _strip_inline_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.replace("`", "")


def rewrite_md_to_pptx(md: str) -> bytes:
    """Convert a rewrite report (## Slide N sections with bullets and
    **Speaker notes:**) back into a .pptx with real notes pages."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    for section in re.split(r"^## ", md, flags=re.M)[1:]:
        lines = section.splitlines()
        heading = _strip_inline_md(lines[0].strip())
        title = re.sub(r"^Slide \d+:\s*", "", heading)
        source = ""
        m = re.search(r"\s*\((source:[^)]*)\)\s*$", title, re.I)
        if m:
            source = m.group(1)
            title = title[: m.start()].strip()

        content: list[str] = []
        notes: list[str] = []
        in_notes = False
        for ln in lines[1:]:
            stripped = ln.strip()
            if re.match(r"\*\*Speaker notes:?\*\*", stripped, re.I):
                in_notes = True
                rest = re.sub(r"^\*\*Speaker notes:?\*\*\s*", "", stripped, flags=re.I)
                if rest:
                    notes.append(rest)
                continue
            if stripped == "---":
                continue
            (notes if in_notes else content).append(ln)

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title or heading
        tf = slide.placeholders[1].text_frame
        first = True
        for ln in content:
            text = ln.rstrip()
            if not text.strip():
                continue
            level = 1 if re.match(r"\s+[-*]\s", ln) else 0
            text = _strip_inline_md(re.sub(r"^\s*(?:[-*]|#+)\s+", "", text)).strip()
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            para.text = text
            para.level = level
            first = False

        notes_text = _strip_inline_md("\n".join(notes).strip())
        if source:
            notes_text = f"[{source}]\n{notes_text}".strip()
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_INLINE_MD_RE = re.compile(r"(\*\*.+?\*\*|`[^`]+`|\*[^*\s][^*]*\*)")


def _add_md_runs(paragraph, text: str) -> None:
    for part in _INLINE_MD_RE.split(text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Courier New"
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            paragraph.add_run(part[1:-1]).italic = True
        else:
            paragraph.add_run(part)


def report_md_to_docx(md: str) -> bytes:
    """Render a markdown report (headings, bullets, tables, blockquotes,
    bold/code) as a .docx document."""
    from docx import Document

    doc = Document()
    lines = md.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped or stripped == "---":
            i += 1
            continue

        if stripped.startswith("|"):  # table block
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                    rows.append(cells)
                i += 1
            if rows:
                table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                table.style = "Table Grid"
                for r, row in enumerate(rows):
                    for c, cell in enumerate(row):
                        para = table.cell(r, c).paragraphs[0]
                        _add_md_runs(para, cell)
                        if r == 0:
                            for run in para.runs:
                                run.bold = True
            continue

        if m := re.match(r"(#{1,6})\s+(.*)", stripped):
            heading = doc.add_heading("", level=min(len(m.group(1)), 4))
            _add_md_runs(heading, m.group(2))
        elif stripped.startswith(">"):
            _add_md_runs(doc.add_paragraph(style="Intense Quote"),
                         stripped.lstrip("> ").strip())
        elif re.match(r"[-*]\s+", stripped):
            style = "List Bullet 2" if lines[i].startswith((" ", "\t")) else "List Bullet"
            _add_md_runs(doc.add_paragraph(style=style),
                         re.sub(r"^[-*]\s+", "", stripped))
        elif re.match(r"\d+\.\s+", stripped):
            _add_md_runs(doc.add_paragraph(style="List Number"),
                         re.sub(r"^\d+\.\s+", "", stripped))
        else:
            _add_md_runs(doc.add_paragraph(), stripped)
        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def export(args) -> int:
    md = Path(args.input).read_text(encoding="utf-8")
    if args.format == "pptx":
        data = rewrite_md_to_pptx(md)
    elif args.format == "docx":
        data = report_md_to_docx(md)
    else:
        print(f"error: unsupported format {args.format}", file=sys.stderr)
        return 2
    Path(args.output).write_bytes(data)
    return 0


# --------------------------------------------------------------------------- #
def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    sub = ap.add_subparsers(dest="cmd", required=True)

    r = sub.add_parser("run", help="run a multi-persona review or rewrite")
    src = r.add_mutually_exclusive_group(required=True)
    src.add_argument("--input", help="path to a document file to extract + review")
    src.add_argument("--text", help="raw pasted text to review")
    r.add_argument("--file-name", default="Document", help="display name for the document")
    r.add_argument("--slug", required=True, help="doc slug (output dir reviews/<slug>/)")
    r.add_argument("--personas", required=True, help="comma-separated persona filenames")
    r.add_argument("--model", required=True)
    r.add_argument("--mode", choices=["review", "rewrite"], default="review")
    r.add_argument("--advise", action="store_true",
                   help="rewrite: draft [DRAFT: …] content for each [NEEDS: …] gap")
    r.add_argument("--context", default="", help="background context text to inform the review/rewrite")
    r.add_argument("--context-file", help="path to a file whose text is used as background context")
    r.add_argument("--ollama-url", default=DEFAULT_OLLAMA_URL)
    r.set_defaults(func=run)

    e = sub.add_parser("export", help="render a report markdown file to docx/pptx")
    e.add_argument("--format", choices=["docx", "pptx"], required=True)
    e.add_argument("--input", required=True, help="path to the report .md file")
    e.add_argument("--output", required=True, help="path to write the exported file")
    e.set_defaults(func=export)

    args = ap.parse_args()
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
