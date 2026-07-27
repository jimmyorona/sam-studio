#!/usr/bin/env python3
"""Extract reviewable text from PPTX, PDF, DOCX, HTML, Markdown, or plain-text files.

Output is markdown on stdout, structured with per-slide / per-page headings so
reviewers can reference locations precisely.

Usage:
    python3 scripts/extract.py <input-file>

Dependencies (all system-level, already present on this machine):
    PPTX -> python-pptx
    PDF  -> pdftotext (poppler-utils)
    DOCX -> libreoffice --headless
    HTML -> stdlib html.parser
"""

import re
import subprocess
import sys
import tempfile
from pathlib import Path


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        out.append(f"## Slide {i}" + (f": {title}" if title else ""))

        body = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text and text != title:
                body.append(text)
        if body:
            out.append("\n\n".join(body))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"**Speaker notes:**\n{notes}")
    return "\n\n".join(out)


def extract_pdf(path: Path) -> str:
    result = subprocess.run(
        ["pdftotext", "-enc", "UTF-8", str(path), "-"],
        capture_output=True, text=True, check=True,
    )
    pages = result.stdout.split("\f")
    out = []
    for i, page in enumerate(pages, 1):
        page = page.strip()
        if page:
            out.append(f"## Page {i}\n\n{page}")
    return "\n\n".join(out)


def extract_docx(path: Path) -> str:
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text",
             "--outdir", tmp, str(path)],
            capture_output=True, check=True,
        )
        txt = Path(tmp) / (path.stem + ".txt")
        return txt.read_text(encoding="utf-8", errors="replace").strip()


def extract_html(path: Path) -> str:
    """HTML slide decks (reveal.js / Marp exports) or plain HTML pages.

    When the document contains multiple <section> elements they are treated as
    slides and numbered, matching the PPTX/PDF heading convention; otherwise the
    page's own heading structure is kept.
    """
    from html.parser import HTMLParser

    raw = path.read_text(encoding="utf-8", errors="replace")
    slide_mode = len(re.findall(r"<section\b", raw, re.IGNORECASE)) >= 2

    class TextExtractor(HTMLParser):
        SKIP = {"script", "style", "head", "template", "noscript", "svg"}
        BLOCK = {"p", "div", "ul", "ol", "table", "tr", "section", "article",
                 "header", "footer", "aside", "blockquote", "pre", "br"}
        HEADINGS = {"h1", "h2", "h3", "h4", "h5", "h6"}

        def __init__(self):
            super().__init__(convert_charrefs=True)
            self.parts = []
            self.skip_depth = 0
            self.slide_no = 0

        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP:
                self.skip_depth += 1
                return
            if self.skip_depth:
                return
            if slide_mode and tag == "section":
                self.slide_no += 1
                self.parts.append(f"\n\n## Slide {self.slide_no}\n\n")
            elif tag in self.HEADINGS:
                # In slide mode, demote headings below the "## Slide N" level.
                level = min(int(tag[1]) + (2 if slide_mode else 0), 6)
                self.parts.append("\n\n" + "#" * level + " ")
            elif tag == "li":
                self.parts.append("\n- ")
            elif tag in self.BLOCK:
                self.parts.append("\n")

        def handle_endtag(self, tag):
            if tag in self.SKIP:
                if self.skip_depth:
                    self.skip_depth -= 1
            elif tag in self.HEADINGS or tag in self.BLOCK or tag == "li":
                self.parts.append("\n")

        def handle_data(self, data):
            if self.skip_depth:
                return
            text = " ".join(data.split())
            if text:
                self.parts.append(text + " ")

    parser = TextExtractor()
    parser.feed(raw)
    parser.close()
    text = "".join(parser.parts)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" ?\n ?", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__, file=sys.stderr)
        return 2
    path = Path(sys.argv[1])
    if not path.exists():
        print(f"error: {path} not found", file=sys.stderr)
        return 1

    ext = path.suffix.lower()
    if ext == ".pptx":
        content = extract_pptx(path)
    elif ext == ".pdf":
        content = extract_pdf(path)
    elif ext in (".docx", ".doc", ".odt"):
        content = extract_docx(path)
    elif ext in (".html", ".htm"):
        content = extract_html(path)
    elif ext in (".md", ".markdown", ".txt"):
        content = path.read_text(encoding="utf-8", errors="replace")
    else:
        print(f"error: unsupported format {ext}", file=sys.stderr)
        return 1

    print(f"# Extracted content: {path.name}\n")
    print(content)
    return 0


if __name__ == "__main__":
    sys.exit(main())
