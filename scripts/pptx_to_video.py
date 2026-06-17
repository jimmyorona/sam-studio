#!/usr/bin/env python3
"""
PPTX / Markdown to AI-Narrated Video Converter

Converts a PowerPoint (.pptx) or Markdown (.md) file into an .mp4 video with
AI-enhanced narration per slide, using a local Ollama model for script generation
and Microsoft Edge TTS (edge-tts) for speech synthesis.

Dependencies:
    pip install edge-tts python-pptx

External tools required:
    - soffice   (LibreOffice)  PPTX → PDF conversion
    - pdftoppm  (poppler)      PDF → PNG extraction
    - ffmpeg                     video assembly
    - npx / @marp-team/marp-cli  Markdown → styled PNG slides

Usage:
    python scripts/pptx_to_video.py \\
        --input ./EA-TECHNICAL-OVERVIEW.pptx \\
        --output outputs/ea-overview.mp4 \\
        --model llama3.2:3b \\
        --voice en-US-AriaNeural \\
        --pause 1.5

    python scripts/pptx_to_video.py \\
        --input ./slides.md \\
        --output outputs/slides.mp4 \\
        --theme gaia
"""

import argparse
import asyncio
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import urllib.request
from pathlib import Path
from typing import List, Optional, Tuple

import urllib.error

try:
    import edge_tts
except ModuleNotFoundError:
    edge_tts = None  # only required when --tts-provider edge (default)

try:
    from supertonic import TTS as SupertonicTTS
except ModuleNotFoundError:
    SupertonicTTS = None  # only required when --tts-provider supertonic

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
DEFAULT_MODEL = "llama3.2:3b"
DEFAULT_VOICE = "en-US-AriaNeural"
DEFAULT_PAUSE = 1.0
DEFAULT_OLLAMA_URL = "http://localhost:11434"
DEFAULT_DPI = 150
DEFAULT_ANIM_FPS = 8
DEFAULT_THEME = "default"
MARP_THEMES = ("default", "gaia", "uncover")

DEFAULT_TTS_PROVIDER = "edge"                        # "edge" | "elevenlabs" | "supertonic"
DEFAULT_ELEVENLABS_VOICE = "pNInz6obpgDQGcFmaJgB"   # Adam — narrator voice
DEFAULT_ELEVENLABS_MODEL = "eleven_multilingual_v2"
ELEVENLABS_API_URL = "https://api.elevenlabs.io/v1/text-to-speech"

# Supertonic — local, on-device ONNX TTS (https://github.com/supertone-inc/supertonic)
DEFAULT_SUPERTONIC_VOICE = "M1"
DEFAULT_SUPERTONIC_LANG = "en"
DEFAULT_SUPERTONIC_STEPS = 8
DEFAULT_SUPERTONIC_SPEED = 1.0

SLIDE_USER_TEMPLATE = (
    "Slide content:\n{slide_text}"
    "{notes_block}"
    "\n\nNarration:"
)

OLLAMA_TIMEOUT = 300  # seconds per slide
FFMPEG_TIMEOUT = 300  # seconds

MUSIC_PROMPT_SYSTEM_BASE = (
    "You are a music director scoring a presentation video. "
    "Write a concise, vivid music description for the ElevenLabs Sound Generation API. "
    "Output plain text only — no markdown, no bullet lists, no stage directions."
)

MUSIC_PROMPT_USER_TEMPLATE = (
    "Presentation slides:\n\n{content}\n\n"
    "Based on the tone of the narrator persona and the content above, "
    "write a single music description (under 120 words) for background music "
    "that fits this entire presentation. Describe genre, mood, tempo, instrumentation, "
    "and any dynamic arc."
)


def log(msg: str) -> None:
    print(f"[{time.strftime('%H:%M:%S')}] {msg}", flush=True)


# ---------------------------------------------------------------------------
# Validation / Prerequisites
# ---------------------------------------------------------------------------
def check_prerequisites() -> None:
    missing = []
    for binary in ("soffice", "pdftoppm", "ffmpeg", "ffprobe"):
        if shutil.which(binary) is None:
            missing.append(binary)
    if missing:
        raise SystemExit(
            f"ERROR: Missing required tools: {', '.join(missing)}. "
            "Please install LibreOffice, poppler-utils (pdftoppm), and FFmpeg."
        )


def find_chrome() -> Optional[str]:
    """Return a path to an existing Chrome/Chromium binary, or None."""
    # Prefer Puppeteer-cached Chrome (already downloaded, version-matched to marp's puppeteer)
    puppeteer_cache = sorted(
        Path.home().glob(".cache/puppeteer/chrome/*/chrome-linux64/chrome"),
        reverse=True,  # newest version first
    )
    candidates: List[Optional[str]] = [
        str(p) for p in puppeteer_cache if p.exists()
    ] + [
        shutil.which("google-chrome-stable"),
        shutil.which("google-chrome"),
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
    ]
    for c in candidates:
        if c and Path(c).exists():
            return c
    return None


def check_marp() -> None:
    """Verify that npx + @marp-team/marp-cli are available."""
    if shutil.which("npx") is None:
        raise SystemExit("ERROR: 'npx' not found. Install Node.js to use Marp for Markdown slides.")
    try:
        result = subprocess.run(
            ["npx", "--yes", "@marp-team/marp-cli", "--version"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip())
        log(f"Marp CLI: {result.stdout.strip()}")
    except Exception as exc:
        raise SystemExit(f"ERROR: Cannot run Marp CLI ({exc}). Run: npm i -g @marp-team/marp-cli")


def check_ollama(ollama_url: str) -> None:
    try:
        req = urllib.request.Request(f"{ollama_url}/api/tags", method="GET")
        with urllib.request.urlopen(req, timeout=10) as resp:
            if resp.status != 200:
                raise RuntimeError(f"Ollama returned status {resp.status}")
    except Exception as exc:
        raise SystemExit(
            f"ERROR: Cannot reach Ollama at {ollama_url} ({exc}). Ensure Ollama is running."
        )


_CTX_CACHE: dict = {}


def model_max_ctx(ollama_url: str, model: str):
    """The selected model's maximum context window, via Ollama /api/show.

    Narration is stateful (the whole deck accumulates in one conversation), so
    using the model's full window avoids silent truncation. Set NARRATE_NUM_CTX
    to override (e.g. to cap memory on a small GPU). Returns None if it can't be
    determined, in which case Ollama's own default applies.
    """
    override = os.environ.get("NARRATE_NUM_CTX")
    if override and override.isdigit():
        return int(override)
    if model in _CTX_CACHE:
        return _CTX_CACHE[model]
    n = None
    try:
        data = json.dumps({"model": model}).encode("utf-8")
        req = urllib.request.Request(
            f"{ollama_url}/api/show", data=data,
            headers={"Content-Type": "application/json"}, method="POST",
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            info = json.loads(resp.read().decode("utf-8")).get("model_info", {}) or {}
        for key, val in info.items():
            if key.endswith(".context_length") and isinstance(val, int):
                n = val
                break
    except Exception:  # best-effort; fall back to Ollama default
        n = None
    _CTX_CACHE[model] = n
    return n


# ---------------------------------------------------------------------------
# Markdown Slide Parsing
# ---------------------------------------------------------------------------
def _strip_markdown_formatting(text: str) -> str:
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    text = re.sub(r'^[-*]\s+', '', text, flags=re.MULTILINE)
    return text.strip()


def parse_markdown_slides(md_path: str) -> List[str]:
    """Split a .md file into per-slide plain-text strings for narration.

    Splits on --- horizontal rules first; if the file has only one block,
    falls back to splitting on level-1 or level-2 headings.
    """
    with open(md_path, 'r', encoding='utf-8') as fh:
        content = fh.read()

    chunks = re.split(r'\n\s*---+\s*\n', content)

    if len(chunks) < 2:
        # Try heading-based split: treat each # / ## as a new slide
        parts = re.split(r'(?=\n#{1,2}\s)', '\n' + content)
        chunks = [p for p in parts if p.strip()]

    slides: List[Tuple[str, str]] = []
    for chunk in chunks:
        is_frame = bool(re.search(r'<!--\s*FRAME\s*-->', chunk, re.IGNORECASE))
        plain = _strip_markdown_formatting(chunk)
        if plain:
            slides.append((plain, "[FRAME]" if is_frame else ""))

    return slides or [("This presentation has no readable text content.", "")]



# ---------------------------------------------------------------------------
# Slide Text Extraction
# ---------------------------------------------------------------------------
def extract_slide_texts(pptx_path: str) -> List[Tuple[str, str]]:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides: List[Tuple[str, str]] = []
    for slide in prs.slides:
        paragraphs: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(("\n".join(paragraphs), notes))
    return slides


# ---------------------------------------------------------------------------
# Slide Grouping (animation frame support)
# ---------------------------------------------------------------------------
def group_slide_data(
    slide_texts: List[Tuple[str, str]],
    image_paths: List[Path],
) -> List[dict]:
    """Group consecutive [FRAME]-tagged slides under their preceding base slide.

    A slide whose speaker notes contain '[FRAME]' is treated as an additional
    animation frame for the previous base slide rather than a new narration beat.

    Returns list of dicts: {"texts": (body, notes), "images": [Path, ...]}
    """
    groups: List[dict] = []
    for (body, notes), img in zip(slide_texts, image_paths):
        if "[FRAME]" in notes.upper() and groups:
            groups[-1]["images"].append(img)
        else:
            groups.append({"texts": (body, notes), "images": [img]})
    return groups


# ---------------------------------------------------------------------------
# Persona Loading
# ---------------------------------------------------------------------------
def load_persona(persona_text: str, persona_file: str) -> str:
    if persona_file:
        path = Path(persona_file)
        if not path.exists():
            raise SystemExit(f"ERROR: Persona file not found: {path}")
        return path.read_text(encoding="utf-8").strip()
    return persona_text.strip()


def load_context(context_text: str, context_file: str) -> str:
    if context_file:
        path = Path(context_file)
        if not path.exists():
            raise SystemExit(f"ERROR: Context file not found: {path}")
        return path.read_text(encoding="utf-8").strip()
    return context_text.strip()


# ---------------------------------------------------------------------------
# Narration Generation (Ollama /api/chat — stateful across slides)
# ---------------------------------------------------------------------------
def build_chat_history(persona: str, context: str) -> List[dict]:
    """Build the initial messages list for a presentation run."""
    parts = [
        "You are narrating a slide presentation. "
        "Write natural spoken narration for each slide you are given. "
        "Output plain text only — no stage directions, sound effects, or markdown formatting."
    ]
    if persona:
        parts.insert(0, persona)
    if context:
        parts.append(
            "Reference document (draw on this as background knowledge when relevant):\n\n" + context
        )
    return [{"role": "system", "content": "\n\n".join(parts)}]


def generate_narration_chat(
    slide_text: str,
    notes_text: str,
    history: List[dict],
    model: str,
    ollama_url: str,
    num_ctx: "int | None" = None,
) -> Tuple[str, List[dict]]:
    """Generate narration for one slide, appending to the full conversation history."""
    if not slide_text.strip():
        narration = "This slide contains visual content with no text to narrate."
        return narration, history + [
            {"role": "user", "content": "Slide content: [visual only — no text]"},
            {"role": "assistant", "content": narration},
        ]

    notes_block = f"\n\nSlide notes:\n{notes_text}" if notes_text.strip() else ""
    user_content = SLIDE_USER_TEMPLATE.format(slide_text=slide_text, notes_block=notes_block)

    messages = history + [{"role": "user", "content": user_content}]
    options = {"temperature": 0.3, "top_p": 0.9}
    if num_ctx:
        options["num_ctx"] = num_ctx
    payload = {
        "model": model,
        "messages": messages,
        "stream": False,
        "options": options,
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        f"{ollama_url}/api/chat",
        data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )

    with urllib.request.urlopen(req, timeout=OLLAMA_TIMEOUT) as resp:
        response = json.loads(resp.read().decode("utf-8"))

    raw = response.get("message", {}).get("content", "").strip()
    raw = re.sub(r"```[\w]*\n?", "", raw)
    raw = raw.replace("```", "").strip()
    narration = raw or " "

    return narration, messages + [{"role": "assistant", "content": narration}]


def generate_music_prompt(
    slide_texts: List[Tuple[str, str]],
    persona: str,
    context: str,
    model: str,
    ollama_url: str,
    num_ctx: "int | None" = None,
) -> str:
    """Generate a single background-music prompt from all slide content."""
    lines = []
    for i, (body, notes) in enumerate(slide_texts, 1):
        lines.append(f"--- Slide {i} ---")
        if body.strip():
            lines.append(body.strip())
        if notes.strip():
            lines.append(f"Notes: {notes.strip()}")
    content = "\n\n".join(lines)

    MAX_CHARS = 6000
    if len(content) > MAX_CHARS:
        content = content[:MAX_CHARS] + "\n\n[...content truncated...]"

    system_parts = [MUSIC_PROMPT_SYSTEM_BASE]
    if persona:
        system_parts.insert(0, persona)
    if context:
        system_parts.append(
            "Reference document (draw on this as background knowledge when relevant):\n\n" + context
        )

    options = {"temperature": 0.6, "top_p": 0.9}
    if num_ctx:
        options["num_ctx"] = num_ctx
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": "\n\n".join(system_parts)},
            {"role": "user", "content": MUSIC_PROMPT_USER_TEMPLATE.format(content=content)},
        ],
        "stream": False,
        "options": options,
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        f"{ollama_url}/api/chat",
        data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )

    with urllib.request.urlopen(req, timeout=OLLAMA_TIMEOUT) as resp:
        response = json.loads(resp.read().decode("utf-8"))

    raw = response.get("message", {}).get("content", "").strip()
    raw = re.sub(r"```[\w]*\n?", "", raw)
    raw = raw.replace("```", "").strip()
    return raw or "Ambient background music, calm and unobtrusive, suitable for a professional presentation."


# ---------------------------------------------------------------------------
# Markdown → styled PNG slides (Marp CLI)
# ---------------------------------------------------------------------------
def render_markdown_with_marp(md_path: str, output_dir: Path, theme: str) -> List[Path]:
    """Render a Markdown slide deck to themed PNG images using Marp CLI.

    Produces output_dir/mslide.001.png, mslide.002.png, ...
    """
    log(f"Rendering Markdown slides with Marp (theme: {theme}) ...")
    # Marp strips the .png suffix from the output base name and appends .001.png, .002.png ...
    # so the prefix must already end in .png: /tmp/work/mslide.png → mslide.001.png, mslide.002.png
    prefix = output_dir / "mslide.png"

    cmd = [
        "npx", "--yes", "@marp-team/marp-cli",
        "--no-stdin",               # don't block waiting for stdin when run as a subprocess
        "--images", "png",
        "--image-scale", "2",       # 2× scale → ~2560×1440 for crisp 1920×1080 video
        "--theme", theme,
        "--allow-local-files",
        "--output", str(prefix),
        md_path,
    ]

    env = os.environ.copy()
    chrome = find_chrome()
    if chrome:
        # Tell Puppeteer (used internally by Marp) to use the already-downloaded Chrome
        # instead of trying to download a new copy, which can time out.
        env["PUPPETEER_EXECUTABLE_PATH"] = chrome
        log(f"Using Chrome: {chrome}")
    else:
        log("WARNING: No local Chrome/Chromium found. Marp may try to download one (~170 MB).")

    result = subprocess.run(cmd, capture_output=True, text=True, timeout=600, env=env)
    if result.returncode != 0:
        raise RuntimeError(f"Marp failed:\n{result.stderr.strip()}")

    pngs = sorted(output_dir.glob("mslide.*.png"))
    if not pngs:
        raise RuntimeError("Marp did not produce any PNG images — check the Markdown file.")
    log(f"Rendered {len(pngs)} slide image(s) with Marp.")
    return pngs


# ---------------------------------------------------------------------------
# Slide Image Rendering (soffice → pdf → pdftoppm → png)
# ---------------------------------------------------------------------------
def render_slides_to_images(pptx_path: str, output_dir: Path, dpi: int) -> List[Path]:
    log("Converting PPTX → PDF with LibreOffice ...")
    stem = Path(pptx_path).stem
    pdf_path = output_dir / f"{stem}.pdf"

    subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            pptx_path,
        ],
        check=True,
        timeout=120,
    )
    if not pdf_path.exists():
        raise RuntimeError(f"LibreOffice did not produce expected PDF: {pdf_path}")
    log(f"PDF created: {pdf_path}")

    log("Extracting slide images with pdftoppm ...")
    prefix = output_dir / "slide"
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
        check=True,
        timeout=120,
    )

    pngs = sorted(output_dir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError("pdftoppm did not produce any PNG images.")
    log(f"Rendered {len(pngs)} slide image(s).")
    return pngs


# ---------------------------------------------------------------------------
# Audio validation helper
# ---------------------------------------------------------------------------
def _validate_audio_file(path: Path, provider: str) -> None:
    if not path.exists():
        raise RuntimeError(f"{provider} did not create audio file: {path}")
    if path.stat().st_size == 0:
        raise RuntimeError(f"{provider} created an empty audio file: {path}")
    # Also verify ffprobe can read it (basic sanity check)
    result = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", str(path)],
        capture_output=True, text=True, timeout=30,
    )
    raw = result.stdout.strip()
    if not raw:
        raise RuntimeError(
            f"{provider} audio file is unreadable by ffprobe: {path}. "
            f"stderr: {result.stderr.strip() or 'None'}"
        )
    try:
        float(raw)
    except ValueError:
        raise RuntimeError(
            f"{provider} audio file has invalid ffprobe duration '{raw}': {path}"
        )


# ---------------------------------------------------------------------------
# Audio Synthesis — Edge TTS (local)
# ---------------------------------------------------------------------------
async def _synthesize_one_edge(narration: str, output_path: Path, voice: str) -> Path:
    if edge_tts is None:
        raise RuntimeError("edge-tts is not installed. Run: pip install edge-tts")
    communicate = edge_tts.Communicate(narration, voice)
    await communicate.save(str(output_path))
    return output_path


async def generate_all_audios_edge(
    narrations: List[str],
    output_dir: Path,
    voice: str,
    concurrency: int = 5,
) -> List[Path]:
    sem = asyncio.Semaphore(concurrency)

    async def _run(idx: int, text: str) -> Tuple[int, Path]:
        out = output_dir / f"slide_{idx + 1:03d}.mp3"
        async with sem:
            await _synthesize_one_edge(text, out, voice)
        _validate_audio_file(out, "Edge TTS")
        return idx, out

    tasks = [_run(i, t) for i, t in enumerate(narrations)]
    log(f"Synthesizing {len(tasks)} audio track(s) with Edge TTS ...")
    results = await asyncio.gather(*tasks)
    paths = [p for _, p in sorted(results, key=lambda x: x[0])]
    for p in paths:
        _validate_audio_file(p, "Edge TTS")
    return paths


# ---------------------------------------------------------------------------
# Audio Synthesis — ElevenLabs (cloud API)
# ---------------------------------------------------------------------------
def _synthesize_one_elevenlabs(
    narration: str,
    output_path: Path,
    *,
    api_key: str,
    voice_id: str,
    model_id: str,
    stability: float,
    similarity_boost: float,
    style: float,
    speed: float,
    use_speaker_boost: bool,
) -> Path:
    payload = json.dumps({
        "text": narration,
        "model_id": model_id,
        "voice_settings": {
            "stability": stability,
            "similarity_boost": similarity_boost,
            "style": style,
            "speed": speed,
            "use_speaker_boost": use_speaker_boost,
        },
    }).encode("utf-8")

    req = urllib.request.Request(
        f"{ELEVENLABS_API_URL}/{voice_id}?output_format=mp3_44100_128",
        data=payload,
        headers={
            "Content-Type": "application/json",
            "xi-api-key": api_key,
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            output_path.write_bytes(resp.read())
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"ElevenLabs API error {exc.code}: {body}") from exc

    return output_path


def generate_all_audios_elevenlabs(
    narrations: List[str],
    output_dir: Path,
    *,
    api_key: str,
    voice_id: str,
    model_id: str,
    stability: float,
    similarity_boost: float,
    style: float,
    speed: float,
    use_speaker_boost: bool,
) -> List[Path]:
    # Synthesize sequentially — ElevenLabs enforces per-minute character limits
    paths: List[Path] = []
    log(f"Synthesizing {len(narrations)} audio track(s) with ElevenLabs ...")
    for idx, text in enumerate(narrations):
        out = output_dir / f"slide_{idx + 1:03d}.mp3"
        log(f"  ElevenLabs: slide {idx + 1}/{len(narrations)} ...")
        _synthesize_one_elevenlabs(
            text, out,
            api_key=api_key,
            voice_id=voice_id,
            model_id=model_id,
            stability=stability,
            similarity_boost=similarity_boost,
            style=style,
            speed=speed,
            use_speaker_boost=use_speaker_boost,
        )
        _validate_audio_file(out, "ElevenLabs")
        paths.append(out)
    for p in paths:
        _validate_audio_file(p, "ElevenLabs")
    return paths


# ---------------------------------------------------------------------------
# Audio Synthesis — Supertonic (local, on-device ONNX)
# ---------------------------------------------------------------------------
def generate_all_audios_supertonic(
    narrations: List[str],
    output_dir: Path,
    *,
    voice: str,
    lang: str,
    steps: int,
    speed: float,
) -> List[Path]:
    if SupertonicTTS is None:
        raise RuntimeError(
            "supertonic is not installed. Run: pip install supertonic"
        )

    # The ONNX model is expensive to load, so build the engine once and reuse
    # it for every slide (and resolve the voice style a single time).
    log("Loading Supertonic model (first run downloads weights) ...")
    tts = SupertonicTTS(auto_download=True)
    style = tts.get_voice_style(voice_name=voice)

    # Supertonic emits 44.1kHz WAV; ffprobe/ffmpeg downstream accept it as-is.
    paths: List[Path] = []
    log(f"Synthesizing {len(narrations)} audio track(s) with Supertonic ...")
    for idx, text in enumerate(narrations):
        out = output_dir / f"slide_{idx + 1:03d}.wav"
        log(f"  Supertonic: slide {idx + 1}/{len(narrations)} ...")
        wav, _duration = tts.synthesize(
            text=text,
            lang=lang,
            voice_style=style,
            total_steps=steps,
            speed=speed,
        )
        tts.save_audio(wav, str(out))
        _validate_audio_file(out, "Supertonic")
        paths.append(out)
    return paths


# ---------------------------------------------------------------------------
# Video Assembly (ffmpeg)
# ---------------------------------------------------------------------------
def get_audio_duration(audio_path: Path) -> float:
    result = subprocess.run(
        [
            "ffprobe",
            "-v", "error",
            "-show_entries", "format=duration",
            "-of", "default=noprint_wrappers=1:nokey=1",
            str(audio_path),
        ],
        capture_output=True,
        text=True,
        timeout=30,
    )
    raw = result.stdout.strip()
    if not raw:
        raise RuntimeError(
            f"ffprobe returned empty duration for {audio_path}. "
            f"stderr: {result.stderr.strip() or 'None'}"
        )
    try:
        return float(raw)
    except ValueError as exc:
        raise RuntimeError(
            f"ffprobe returned non-numeric duration '{raw}' for {audio_path}"
        ) from exc


def build_concat_list(
    image_groups: List[List[Path]],
    audio_paths: List[Path],
    work_dir: Path,
    pause: float,
    anim_fps: int = DEFAULT_ANIM_FPS,
) -> Tuple[Path, List[Path]]:
    concat_file = work_dir / "concat_list.txt"
    segment_paths: List[Path] = []
    vf = (
        "scale=1920:1080:force_original_aspect_ratio=decrease,"
        "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:black"
    )

    for idx, (imgs, aud) in enumerate(zip(image_groups, audio_paths)):
        duration = get_audio_duration(aud) + pause
        segment = work_dir / f"segment_{idx + 1:03d}.mp4"

        if len(imgs) == 1:
            subprocess.run(
                [
                    "ffmpeg", "-y",
                    "-loop", "1", "-i", str(imgs[0]),
                    "-i", str(aud),
                    "-vf", vf,
                    "-c:v", "libx264", "-tune", "stillimage",
                    "-c:a", "aac", "-b:a", "192k",
                    "-pix_fmt", "yuv420p",
                    "-t", str(duration), "-r", "30", "-shortest",
                    str(segment),
                ],
                check=True,
                timeout=FFMPEG_TIMEOUT,
            )
        else:
            # Animated group: loop frames at anim_fps over the full audio duration
            frame_dur = 1.0 / anim_fps
            loop_count = math.ceil(duration / (len(imgs) * frame_dur)) + 1
            frames_file = work_dir / f"frames_{idx + 1:03d}.txt"
            with open(frames_file, "w") as fh:
                for _ in range(loop_count):
                    for img in imgs:
                        fh.write(f"file '{img}'\nduration {frame_dur:.6f}\n")
                fh.write(f"file '{imgs[-1]}'\n")  # ffmpeg concat requires a trailing entry
            subprocess.run(
                [
                    "ffmpeg", "-y",
                    "-f", "concat", "-safe", "0", "-i", str(frames_file),
                    "-i", str(aud),
                    "-vf", vf,
                    "-c:v", "libx264",
                    "-c:a", "aac", "-b:a", "192k",
                    "-pix_fmt", "yuv420p",
                    "-r", "30", "-shortest",
                    str(segment),
                ],
                check=True,
                timeout=FFMPEG_TIMEOUT,
            )

        segment_paths.append(segment)

    with open(concat_file, "w") as fh:
        for seg in segment_paths:
            fh.write(f"file '{seg}'\n")

    return concat_file, segment_paths


def concatenate_segments(concat_file: Path, output_path: Path) -> None:
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-f", "concat",
            "-safe", "0",
            "-i", str(concat_file),
            "-c", "copy",
            str(output_path),
        ],
        check=True,
        timeout=FFMPEG_TIMEOUT,
    )


# ---------------------------------------------------------------------------
# Main entrypoint
# ---------------------------------------------------------------------------
def main() -> int:
    parser = argparse.ArgumentParser(description="Convert a PPTX or Markdown file to an AI-narrated MP4 video.")
    parser.add_argument("--input", "-i", required=True, help="Path to input .pptx or .md file")
    parser.add_argument("--output", "-o", required=True, help="Path for output .mp4 file")
    parser.add_argument("--model", "-m", default=DEFAULT_MODEL, help=f"Ollama model (default: {DEFAULT_MODEL})")
    parser.add_argument("--pause", "-p", type=float, default=DEFAULT_PAUSE, help=f"Pause after each slide (default: {DEFAULT_PAUSE}s)")
    parser.add_argument("--dpi", "-d", type=int, default=DEFAULT_DPI, help=f"PNG DPI for PPTX rendering (default: {DEFAULT_DPI})")
    parser.add_argument("--anim-fps", type=int, default=DEFAULT_ANIM_FPS,
                        help=f"Frames per second for [FRAME]-tagged animation slide groups (default: {DEFAULT_ANIM_FPS})")
    parser.add_argument("--theme", "-t", default=DEFAULT_THEME, choices=list(MARP_THEMES),
                        help=f"Marp theme for Markdown slides (default: {DEFAULT_THEME})")
    parser.add_argument("--ollama-url", "-u", default=DEFAULT_OLLAMA_URL, help=f"Ollama base URL (default: {DEFAULT_OLLAMA_URL})")
    parser.add_argument("--keep-temp", action="store_true", help="Preserve temporary working directory")
    parser.add_argument(
        "--manifest-output", default="",
        help="Write a JSON manifest of slide images and narrations to this path and exit before TTS/assembly",
    )
    parser.add_argument(
        "--generate-music-prompt-only",
        action="store_true",
        help="Extract slide text, ask Ollama for a music prompt, write JSON to --output, and exit.",
    )

    # Persona / prompt injection
    persona_group = parser.add_argument_group("Narrator Persona")
    persona_group.add_argument(
        "--persona-text", default="",
        help="Character brief / tone guide prepended to every narration prompt (raw text)",
    )
    persona_group.add_argument(
        "--persona-file", default="",
        help="Path to a .md or .txt file whose content is used as the narrator persona",
    )
    persona_group.add_argument(
        "--context-text", default="",
        help="Background / reference text injected into every slide prompt (raw text)",
    )
    persona_group.add_argument(
        "--context-file", default="",
        help="Path to a background document (.md or .txt) injected into every slide prompt as reference material",
    )

    # TTS provider
    parser.add_argument(
        "--tts-provider", default=DEFAULT_TTS_PROVIDER,
        choices=["edge", "elevenlabs", "supertonic"],
        help="TTS engine: 'edge' (cloud, free), 'elevenlabs' (cloud API), or "
             "'supertonic' (local, on-device, offline). Default: edge",
    )

    # Edge TTS
    edge_group = parser.add_argument_group("Edge TTS (--tts-provider edge)")
    edge_group.add_argument(
        "--voice", "-v", default=DEFAULT_VOICE,
        help=f"Edge-TTS voice name (default: {DEFAULT_VOICE})",
    )

    # ElevenLabs
    el_group = parser.add_argument_group("ElevenLabs TTS (--tts-provider elevenlabs)")
    el_group.add_argument(
        "--elevenlabs-api-key", default=os.environ.get("ELEVENLABS_API_KEY", ""),
        help="ElevenLabs API key (or set ELEVENLABS_API_KEY env var)",
    )
    el_group.add_argument(
        "--elevenlabs-voice-id", default=DEFAULT_ELEVENLABS_VOICE,
        help=f"ElevenLabs voice ID (default: {DEFAULT_ELEVENLABS_VOICE} — Adam)",
    )
    el_group.add_argument(
        "--elevenlabs-model", default=DEFAULT_ELEVENLABS_MODEL,
        help=f"ElevenLabs model ID (default: {DEFAULT_ELEVENLABS_MODEL})",
    )
    el_group.add_argument(
        "--elevenlabs-stability", type=float, default=0.5,
        help="Voice stability 0–1 (default: 0.5)",
    )
    el_group.add_argument(
        "--elevenlabs-similarity", type=float, default=0.75,
        help="Similarity boost 0–1 (default: 0.75)",
    )
    el_group.add_argument(
        "--elevenlabs-style", type=float, default=0.0,
        help="Style exaggeration 0–1 (default: 0.0)",
    )
    el_group.add_argument(
        "--elevenlabs-speed", type=float, default=1.0,
        help="Speaking speed 0.7–1.2 (default: 1.0)",
    )
    el_group.add_argument(
        "--elevenlabs-no-speaker-boost", action="store_true",
        help="Disable speaker boost (enabled by default)",
    )

    # Supertonic
    st_group = parser.add_argument_group("Supertonic TTS (--tts-provider supertonic)")
    st_group.add_argument(
        "--supertonic-voice", default=DEFAULT_SUPERTONIC_VOICE,
        help=f"Supertonic voice name (default: {DEFAULT_SUPERTONIC_VOICE})",
    )
    st_group.add_argument(
        "--supertonic-lang", default=DEFAULT_SUPERTONIC_LANG,
        help=f"Supertonic language code (default: {DEFAULT_SUPERTONIC_LANG})",
    )
    st_group.add_argument(
        "--supertonic-steps", type=int, default=DEFAULT_SUPERTONIC_STEPS,
        help=f"Diffusion steps — higher is slower/better (default: {DEFAULT_SUPERTONIC_STEPS})",
    )
    st_group.add_argument(
        "--supertonic-speed", type=float, default=DEFAULT_SUPERTONIC_SPEED,
        help=f"Speaking speed (default: {DEFAULT_SUPERTONIC_SPEED})",
    )

    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    output_path = Path(args.output).resolve()

    if not input_path.exists():
        raise SystemExit(f"ERROR: Input file not found: {input_path}")

    suffix = input_path.suffix.lower()
    if suffix not in (".pptx", ".md"):
        raise SystemExit(f"ERROR: Input must be a .pptx or .md file: {input_path}")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    log("Step 1/5: Checking prerequisites ...")
    if not args.generate_music_prompt_only:
        check_prerequisites()
    check_ollama(args.ollama_url)

    # Use the model's full context window so the stateful narration history
    # (which grows with every slide) isn't silently truncated by Ollama.
    num_ctx = model_max_ctx(args.ollama_url, args.model)
    if num_ctx:
        log(f"Model context window: {num_ctx} tokens.")

    work_dir = Path(tempfile.mkdtemp(prefix="pptx_video_"))
    try:
        log(f"Working directory: {work_dir}")

        if suffix == ".md":
            log("Step 2/5: Parsing Markdown slides ...")
            slide_texts = parse_markdown_slides(str(input_path))
            log(f"Found {len(slide_texts)} slide(s) in Markdown.")
            if not args.generate_music_prompt_only and all(not body.strip() for body, _ in slide_texts):
                raise SystemExit("ERROR: All slides are empty; nothing to narrate.")

            if args.generate_music_prompt_only:
                persona = load_persona(args.persona_text, args.persona_file)
                context = load_context(args.context_text, args.context_file)
                prompt = generate_music_prompt(slide_texts, persona, context, args.model, args.ollama_url, num_ctx)
                Path(args.output).write_text(json.dumps({"prompt": prompt}, indent=2), encoding="utf-8")
                log(f"Music prompt written to: {args.output}")
                return 0

            log("Step 3/5: Rendering themed slides with Marp ...")
            check_marp()
            image_paths = render_markdown_with_marp(str(input_path), work_dir, args.theme)
        else:
            log("Step 2/5: Extracting slide text via python-pptx ...")
            slide_texts = extract_slide_texts(str(input_path))
            log(f"Found {len(slide_texts)} slide(s).")
            if not args.generate_music_prompt_only and all(not body.strip() for body, _ in slide_texts):
                raise SystemExit("ERROR: All slides are empty; nothing to narrate.")

            if args.generate_music_prompt_only:
                persona = load_persona(args.persona_text, args.persona_file)
                context = load_context(args.context_text, args.context_file)
                prompt = generate_music_prompt(slide_texts, persona, context, args.model, args.ollama_url, num_ctx)
                Path(args.output).write_text(json.dumps({"prompt": prompt}, indent=2), encoding="utf-8")
                log(f"Music prompt written to: {args.output}")
                return 0

            log("Step 3/5: Rendering slides via LibreOffice ...")
            image_paths = render_slides_to_images(str(input_path), work_dir, args.dpi)
        effective = min(len(image_paths), len(slide_texts))
        log(f"Effective slides to process: {effective}")
        groups = group_slide_data(slide_texts[:effective], image_paths[:effective])
        if len(groups) < effective:
            log(f"Animation groups: {len(groups)} (from {effective} slides; {effective - len(groups)} frame slide(s))")

        persona = load_persona(args.persona_text, args.persona_file)
        if persona:
            log(f"Persona loaded: {persona[:80].splitlines()[0]} ...")

        context = load_context(args.context_text, args.context_file)
        if context:
            src = args.context_file or "inline text"
            log(f"Context loaded: {len(context)} chars from {src}")

        log("Step 4/5: Generating narrations via Ollama ...")
        narrations: List[str] = []
        history = build_chat_history(persona, context)
        for idx, group in enumerate(groups):
            stext, notes = group["texts"]
            notes_for_prompt = re.sub(r'\[FRAME\]', '', notes, flags=re.IGNORECASE).strip()
            narration, history = generate_narration_chat(stext, notes_for_prompt, history, args.model, args.ollama_url, num_ctx)
            narrations.append(narration)
            log(f"  Group {idx + 1}/{len(groups)}: {narration[:100]} ...")

        if args.manifest_output:
            manifest = {
                "work_dir": str(work_dir),
                "images": [[str(p) for p in g["images"]] for g in groups],
                "narrations": narrations,
                "pause": args.pause,
                "anim_fps": args.anim_fps,
            }
            Path(args.manifest_output).write_text(json.dumps(manifest, indent=2), encoding="utf-8")
            log(f"Manifest written: {args.manifest_output} ({len(groups)} group(s) from {effective} slides)")
            return 0

        log("Step 5/5: Synthesizing audio and assembling video ...")
        if args.tts_provider == "elevenlabs":
            if not args.elevenlabs_api_key:
                raise SystemExit(
                    "ERROR: --elevenlabs-api-key is required when using ElevenLabs TTS "
                    "(or set the ELEVENLABS_API_KEY environment variable)."
                )
            audio_paths = generate_all_audios_elevenlabs(
                narrations, work_dir,
                api_key=args.elevenlabs_api_key,
                voice_id=args.elevenlabs_voice_id,
                model_id=args.elevenlabs_model,
                stability=args.elevenlabs_stability,
                similarity_boost=args.elevenlabs_similarity,
                style=args.elevenlabs_style,
                speed=args.elevenlabs_speed,
                use_speaker_boost=not args.elevenlabs_no_speaker_boost,
            )
        elif args.tts_provider == "supertonic":
            audio_paths = generate_all_audios_supertonic(
                narrations, work_dir,
                voice=args.supertonic_voice,
                lang=args.supertonic_lang,
                steps=args.supertonic_steps,
                speed=args.supertonic_speed,
            )
        else:
            audio_paths = asyncio.run(generate_all_audios_edge(narrations, work_dir, args.voice))

        image_groups = [g["images"] for g in groups]
        concat_file, _ = build_concat_list(image_groups, audio_paths, work_dir, args.pause, args.anim_fps)
        concatenate_segments(concat_file, output_path)

        size_mb = output_path.stat().st_size / (1024 * 1024)
        log(f"Done! Output written to: {output_path} ({size_mb:.1f} MB)")
        return 0

    except Exception as exc:
        log(f"Error during processing: {exc}")
        return 1

    finally:
        if args.keep_temp or args.manifest_output:
            log(f"Temp directory kept at: {work_dir}")
        else:
            shutil.rmtree(work_dir, ignore_errors=True)
            log("Temp directory cleaned up.")


if __name__ == "__main__":
    sys.exit(main())
